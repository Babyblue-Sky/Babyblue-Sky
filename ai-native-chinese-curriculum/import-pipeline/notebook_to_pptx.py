"""
Import Pipeline: Extractor stage (Blueprint v1.0, section 7) — SMART Notebook -> Slides.

Reads a .notebook file (a zip: per-page SVGs with positioned <text>/<image>
elements, plus imsmanifest.xml recording the real page order) and rebuilds
it as a .pptx deck, landscape, ready to open in Google Slides / PowerPoint.

This is NOT the same Extractor path that feeds the Canonical Content Layer
(that one still goes Extractor -> AI Classifier -> Human Review -> Markdown,
per the Import Pipeline diagram in blueprint-v1.0.md, and stays a curated,
human-reviewed distillation). This script is a sibling use of the same
underlying SVG-extraction technique, for a different purpose: the teacher
decided (2026-08) to stop authoring/teaching from SMART Notebook and move
her live classroom presentations to Google Slides, decoupled from the
Content Layer's git-driven update cadence (see PROJECT_STATUS.md — live
teaching material needs to be instantly editable, which the Content
Layer/Generation Layer pipeline is deliberately NOT). So this tool
reconstructs a close visual match of the original SMART pages as an
editable Slides deck, not a curated Markdown update.

What it preserves: text (position, size, color) and images (position,
size, transparency), page order from imsmanifest.xml, per-page background
color. What it does NOT preserve: SMART-specific interactive widgets, ink
annotations recorded only as annotationmetadata, animated GIF frames
(first frame only), or exact original image resolution (images are
resized down — see `_prepare_images` — since full-resolution originals
make the .pptx too large to hand off in one piece; that's a deliberate
trade-off for a usable draft, not a bug).

Known limitations worth knowing before you rely on this for a real class:
- Original SMART pages are near-square (~800x700-850px). Standard Slides
  widescreen (960x540pt, 16:9) is much wider. Each page is scaled to fit
  and centered, which leaves visible side margins on most pages (filled
  with that page's own background color, so it reads as generous
  whitespace rather than empty bars — but it's still whitespace, not a
  redesigned wide layout).
- Chinese text is tagged with font "Kaiti SC" (a macOS system font). It
  renders correctly in apps with local access to that font (Keynote,
  PowerPoint on a Mac that has it installed), but Google Slides runs in
  the browser and can't reach local system fonts, so Chinese text shows
  in Slides' fallback font until manually reassigned there (once, e.g. at
  the Slide Master level — not a bug in this script, a cross-platform font
  distribution constraint. See 08-curriculum-intelligence.md.).

Usage:
    python3 notebook_to_pptx.py <path/to/file.notebook> <output.pptx> [--title "Deck title"]
"""
import argparse
import json
import os
import re
import shutil
import sys
import tempfile
import xml.etree.ElementTree as ET
import zipfile

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

SLIDE_W, SLIDE_H = 960.0, 540.0  # pt, standard 13.33x7.5in widescreen
CJK_RE = re.compile(r'[一-鿿㐀-䶿]')


def unzip_notebook(notebook_path, dest_dir):
    with zipfile.ZipFile(notebook_path) as zf:
        zf.extractall(dest_dir)
    return dest_dir


def page_order(extract_dir):
    manifest_path = os.path.join(extract_dir, "imsmanifest.xml")
    manifest = open(manifest_path, encoding="utf-8").read()
    m = re.search(r'<resource identifier="group0_pages"[^>]*>(.*?)</resource>', manifest, re.S)
    files = re.findall(r'<file href="(page[\w.]+\.svg)"/>', m.group(1))
    seen, ordered = set(), []
    for f in files:
        if f not in seen:
            seen.add(f)
            ordered.append(f)
    return ordered


def _parse_translate(t):
    m = re.search(r'translate\(([-\d.]+),\s*([-\d.]+)\)', t or "")
    return (float(m.group(1)), float(m.group(2))) if m else (0.0, 0.0)


def extract_page(svg_path):
    tree = ET.parse(svg_path)
    root = tree.getroot()
    w, h = float(root.get("width")), float(root.get("height"))
    bg, items = None, []
    for el in root.iter():
        tag = el.tag.split("}")[-1]
        if tag == "rect" and bg is None and el.get("fill") and el.get("width") == "100%":
            bg = el.get("fill")
        elif tag == "image":
            items.append({"type": "image", "x": float(el.get("x", 0)), "y": float(el.get("y", 0)),
                          "w": float(el.get("width", 0)), "h": float(el.get("height", 0)),
                          "href": el.get("{http://www.w3.org/1999/xlink}href")})
        elif tag == "text":
            tx, ty = _parse_translate(el.get("transform"))
            runs, max_font, color = [], 0, "#000000"
            for tspan in el.iter():
                if tspan.tag.split("}")[-1] == "tspan" and tspan.text:
                    fs = tspan.get("font-size")
                    if fs:
                        max_font = max(max_font, float(fs))
                    if tspan.get("fill"):
                        color = tspan.get("fill")
                    runs.append(tspan.text)
            if runs:
                items.append({"type": "text", "x": tx, "y": ty, "font_size": max_font or 18,
                              "color": color, "text": "".join(runs)})
    return {"w": w, "h": h, "bg": bg, "items": items}


def _has_real_alpha(im):
    if im.mode not in ("RGBA", "LA", "P"):
        return False
    alpha = im.convert("RGBA").split()[3]
    return alpha.getextrema()[0] < 250


def prepare_images(extract_dir, pages_data, out_dir):
    """Resize every image used across all pages to ~1.5x its largest
    on-page display size, preserving transparency where the source has it
    (JPEG otherwise, since it compresses far better for opaque photos)."""
    os.makedirs(out_dir, exist_ok=True)
    display_size = {}
    for data in pages_data:
        scale = min(SLIDE_W / data["w"], SLIDE_H / data["h"])
        for it in data["items"]:
            if it["type"] != "image":
                continue
            href = it["href"]
            dw, dh = it["w"] * scale, it["h"] * scale
            cur = display_size.get(href, (0, 0))
            display_size[href] = (max(cur[0], dw), max(cur[1], dh))

    mapping = {}
    for href, (dw, dh) in display_size.items():
        src = os.path.join(extract_dir, href)
        if not os.path.exists(src):
            continue
        try:
            im = Image.open(src)
        except Exception as e:
            print(f"  [skip unreadable image {href}: {e}]")
            continue
        if getattr(im, "is_animated", False):
            im.seek(0)
        target_w = max(min(im.width, int(dw * 1.5) or 1), 1)
        target_h = max(min(im.height, int(dh * 1.5) or 1), 1)
        transparent = _has_real_alpha(im)
        im = im.convert("RGBA")
        im.thumbnail((target_w, target_h), Image.LANCZOS)

        base = os.path.splitext(os.path.basename(href))[0]
        # href may repeat a basename across subfolders; keep it unique
        safe = re.sub(r'[^A-Za-z0-9_.-]', "_", base) + f"_{abs(hash(href)) % 100000}"
        if transparent:
            out_name = safe + ".png"
            im.save(os.path.join(out_dir, out_name), "PNG", optimize=True)
        else:
            bg = Image.new("RGB", im.size, (255, 255, 255))
            bg.paste(im, mask=im.split()[3])
            out_name = safe + ".jpg"
            bg.save(os.path.join(out_dir, out_name), "JPEG", quality=75, optimize=True)
        mapping[href] = out_name
    return mapping


def _hex_to_rgb(h):
    h = (h or "#000000").lstrip("#")
    if len(h) != 6:
        return RGBColor(0, 0, 0)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def build_pptx(extract_dir, page_files, image_dir, image_map, out_path, title=None):
    prs = Presentation()
    prs.slide_width = Pt(SLIDE_W)
    prs.slide_height = Pt(SLIDE_H)
    blank = prs.slide_layouts[6]

    if title:
        cover = prs.slides.add_slide(blank)
        tb = cover.shapes.add_textbox(Pt(60), Pt(220), Pt(SLIDE_W - 120), Pt(140))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.name = "Kaiti SC"

    for page_file in page_files:
        data = extract_page(os.path.join(extract_dir, page_file))
        page_w, page_h = data["w"], data["h"]
        scale = min(SLIDE_W / page_w, SLIDE_H / page_h)
        dx = (SLIDE_W - page_w * scale) / 2
        dy = (SLIDE_H - page_h * scale) / 2

        slide = prs.slides.add_slide(blank)
        if data["bg"]:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = _hex_to_rgb(data["bg"])

        for it in data["items"]:
            x = it["x"] * scale + dx
            y = it["y"] * scale + dy
            if it["type"] == "image":
                mapped = image_map.get(it["href"])
                if not mapped:
                    continue
                img_path = os.path.join(image_dir, mapped)
                w, h = it["w"] * scale, it["h"] * scale
                if os.path.exists(img_path) and w > 0 and h > 0:
                    try:
                        slide.shapes.add_picture(img_path, Pt(x), Pt(y), Pt(w), Pt(h))
                    except Exception as e:
                        print(f"  [skip image {it['href']}: {e}]")
            else:
                font_size = it["font_size"] * scale * 0.75
                width = max(Pt((SLIDE_W - x) - 5), Pt(40))
                height = Pt(font_size * 1.6 + 8)
                tb = slide.shapes.add_textbox(Pt(x), Pt(y), width, height)
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = it["text"]
                p.font.size = Pt(max(font_size, 6))
                p.font.color.rgb = _hex_to_rgb(it["color"])
                if CJK_RE.search(it["text"]):
                    p.font.name = "Kaiti SC"

    prs.save(out_path)
    return out_path


def convert(notebook_path, out_path, title=None, keep_temp=False):
    tmp_dir = tempfile.mkdtemp(prefix="notebook_extract_")
    try:
        unzip_notebook(notebook_path, tmp_dir)
        pages = page_order(tmp_dir)
        pages_data = [extract_page(os.path.join(tmp_dir, p)) for p in pages]
        image_dir = os.path.join(tmp_dir, "_prepared_images")
        image_map = prepare_images(tmp_dir, pages_data, image_dir)
        build_pptx(tmp_dir, pages, image_dir, image_map, out_path, title=title)
        print(f"wrote {out_path} ({len(pages)} pages) — {os.path.getsize(out_path)} bytes")
    finally:
        if not keep_temp:
            shutil.rmtree(tmp_dir, ignore_errors=True)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("notebook_path")
    parser.add_argument("out_path")
    parser.add_argument("--title", default=None)
    parser.add_argument("--keep-temp", action="store_true")
    args = parser.parse_args()
    convert(args.notebook_path, args.out_path, title=args.title, keep_temp=args.keep_temp)
